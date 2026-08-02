#!/usr/bin/env python3
"""PowerPoint presentation creator using python-pptx."""

import argparse
from pathlib import Path


def _ensure_deps():
    try:
        from echo_agent.dependencies.skill_require import require
        require("skill.ppt-author")
    except ImportError:
        pass
    require("skill.ppt-author")


def create_pptx(title, slides_data, output="output.pptx", template=None):
    from pptx import Presentation
    from pptx.util import Inches, Pt  # noqa: F401

    prs = Presentation(template) if template else Presentation()

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = title
    if title_slide.placeholders[1]:
        title_slide.placeholders[1].text = ""

    for slide_info in slides_data:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = slide_info.get("title", "")
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.text = ""
        for i, bullet in enumerate(slide_info.get("bullets", [])):
            if i == 0:
                tf.text = bullet
            else:
                p = tf.add_paragraph()
                p.text = bullet
                p.level = slide_info.get("level", 0)

    prs.save(output)
    print(f"Presentation saved: {output} ({len(slides_data) + 1} slides)")


def from_markdown(md_file, output="output.pptx"):
    text = Path(md_file).read_text()
    lines = text.strip().split("\n")
    title = lines[0].lstrip("# ").strip() if lines else "Untitled"
    slides = []
    current = None
    for line in lines[1:]:
        if line.startswith("## "):
            if current:
                slides.append(current)
            current = {"title": line[3:].strip(), "bullets": []}
        elif line.startswith("- ") and current:
            current["bullets"].append(line[2:].strip())
        elif line.strip() and current and current["bullets"]:
            current["bullets"][-1] += " " + line.strip()
    if current:
        slides.append(current)
    create_pptx(title, slides, output)


def main():
    parser = argparse.ArgumentParser(description="PPTX creator")
    sub = parser.add_subparsers(dest="cmd")
    p = sub.add_parser("from-md")
    p.add_argument("input", help="Markdown file")
    p.add_argument("--output", "-o", default="output.pptx")
    p = sub.add_parser("quick")
    p.add_argument("title")
    p.add_argument("--slides", nargs="+", help="Slide titles")
    p.add_argument("--output", "-o", default="output.pptx")
    args = parser.parse_args()

    if args.cmd == "from-md":
        _ensure_deps()
        from_markdown(args.input, args.output)
    elif args.cmd == "quick":
        _ensure_deps()
        slides = [{"title": s, "bullets": []} for s in (args.slides or [])]
        create_pptx(args.title, slides, args.output)
    else:
        parser.print_help()


if __name__ == "__main__":
    main()
