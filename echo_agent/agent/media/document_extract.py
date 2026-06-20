"""Document text extraction core — pure local-path → text, no network/decrypt/model.

Shared by inbound auto-injection (agent/context.py) and the read_document tool.
Every branch degrades safely: missing deps, corrupt files, and unsupported
formats return an empty ExtractResult with a meta marker, never raise."""
from __future__ import annotations

from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

from loguru import logger

_TEXT_EXTS = {".txt", ".md", ".csv", ".json", ".yaml", ".yml", ".log", ".xml", ".html"}


@dataclass
class ExtractResult:
    text: str = ""
    meta: dict[str, Any] = field(default_factory=dict)
    truncated: bool = False
    unit_count: int = 0


def _apply_cap(text: str, max_chars: int | None) -> tuple[str, bool]:
    if max_chars is not None and max_chars >= 0 and len(text) > max_chars:
        return text[:max_chars], True
    return text, False


def extract(
    path: str | Path,
    *,
    max_chars: int | None = None,
    unit: int | str | None = None,
) -> ExtractResult:
    p = Path(path)
    if not p.exists() or not p.is_file():
        return ExtractResult(meta={"format": "missing", "error": f"file not found: {p}"})
    ext = p.suffix.lower()
    size = p.stat().st_size
    try:
        if ext == ".pdf":
            return _extract_pdf(p, max_chars, unit, size)
        if ext == ".docx":
            return _extract_docx(p, max_chars, size)
        if ext == ".xlsx":
            return _extract_xlsx(p, max_chars, unit, size)
        if ext == ".pptx":
            return _extract_pptx(p, max_chars, unit, size)
        if ext in _TEXT_EXTS:
            return _extract_text(p, max_chars, size)
    except Exception as exc:  # noqa: BLE001 — inbound path must never raise
        logger.warning("document extract failed for {}: {}", p.name, exc)
        return ExtractResult(meta={"format": ext.lstrip("."), "size_bytes": size, "error": str(exc)})
    return ExtractResult(meta={"format": "unsupported", "size_bytes": size})


def _extract_text(p: Path, max_chars: int | None, size: int) -> ExtractResult:
    raw = p.read_text(encoding="utf-8", errors="replace")
    text, truncated = _apply_cap(raw, max_chars)
    return ExtractResult(
        text=text,
        meta={"format": p.suffix.lower().lstrip("."), "size_bytes": size},
        truncated=truncated,
        unit_count=raw.count("\n") + 1,
    )


def _extract_pdf(p: Path, max_chars: int | None, unit: int | str | None, size: int) -> ExtractResult:
    try:
        import pymupdf  # type: ignore
    except ImportError:
        return ExtractResult(meta={"format": "pdf", "size_bytes": size,
                                   "error": "missing_dep:pymupdf>=1.24"})
    doc = pymupdf.open(str(p))
    try:
        total = doc.page_count
        parts: list[str] = []
        target_page = int(unit) if isinstance(unit, (int, str)) and str(unit).isdigit() else None
        for i in range(total):
            if target_page is not None and (i + 1) != target_page:
                continue
            page_text = doc[i].get_text()
            parts.append(f"--- 第{i + 1}页 ---\n{page_text}")
            if max_chars is not None and sum(len(x) for x in parts) >= max_chars:
                break
        joined = "\n".join(parts)
        text, truncated = _apply_cap(joined, max_chars)
        return ExtractResult(text=text, meta={"format": "pdf", "size_bytes": size, "pages": total},
                             truncated=truncated, unit_count=total)
    finally:
        doc.close()


def _extract_docx(p: Path, max_chars: int | None, size: int) -> ExtractResult:
    try:
        import docx  # type: ignore
    except ImportError:
        return ExtractResult(meta={"format": "docx", "size_bytes": size,
                                   "error": "missing_dep:python-docx>=1.1"})
    doc = docx.Document(str(p))
    lines = [para.text for para in doc.paragraphs if para.text]
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text for c in row.cells if c.text]
            if cells:
                lines.append("\t".join(cells))
    text, truncated = _apply_cap("\n".join(lines), max_chars)
    return ExtractResult(text=text, meta={"format": "docx", "size_bytes": size},
                         truncated=truncated, unit_count=len(doc.paragraphs))


def _extract_xlsx(p: Path, max_chars: int | None, unit: int | str | None, size: int) -> ExtractResult:
    try:
        import openpyxl  # type: ignore
    except ImportError:
        return ExtractResult(meta={"format": "xlsx", "size_bytes": size,
                                   "error": "missing_dep:openpyxl>=3.1"})
    wb = openpyxl.load_workbook(str(p), read_only=True, data_only=True)
    try:
        names = wb.sheetnames
        if unit is not None:
            if isinstance(unit, int) or (isinstance(unit, str) and unit.isdigit()):
                idx = int(unit) - 1
                targets = [names[idx]] if 0 <= idx < len(names) else []
            else:
                targets = [unit] if unit in names else []
        else:
            targets = names
        parts: list[str] = []
        for name in targets:
            ws = wb[name]
            rows: list[str] = []
            for row in ws.iter_rows(values_only=True):
                rows.append(",".join("" if v is None else str(v) for v in row))
            parts.append(f"=== sheet: {name} ===\n" + "\n".join(rows))
            if max_chars is not None and sum(len(x) for x in parts) >= max_chars:
                break
        text, truncated = _apply_cap("\n\n".join(parts), max_chars)
        return ExtractResult(text=text, meta={"format": "xlsx", "size_bytes": size, "sheets": names},
                             truncated=truncated, unit_count=len(names))
    finally:
        wb.close()


def _extract_pptx(p: Path, max_chars: int | None, unit: int | str | None, size: int) -> ExtractResult:
    try:
        import pptx  # type: ignore
    except ImportError:
        return ExtractResult(meta={"format": "pptx", "size_bytes": size,
                                   "error": "missing_dep:python-pptx>=1.0"})
    prs = pptx.Presentation(str(p))
    slides = list(prs.slides)
    target = int(unit) if isinstance(unit, (int, str)) and str(unit).isdigit() else None
    parts: list[str] = []
    for i, slide in enumerate(slides):
        if target is not None and (i + 1) != target:
            continue
        texts = [sh.text_frame.text for sh in slide.shapes
                 if sh.has_text_frame and sh.text_frame.text]
        parts.append(f"--- 第{i + 1}页 ---\n" + "\n".join(texts))
        if max_chars is not None and sum(len(x) for x in parts) >= max_chars:
            break
    text, truncated = _apply_cap("\n".join(parts), max_chars)
    return ExtractResult(text=text, meta={"format": "pptx", "size_bytes": size, "slides": len(slides)},
                         truncated=truncated, unit_count=len(slides))
