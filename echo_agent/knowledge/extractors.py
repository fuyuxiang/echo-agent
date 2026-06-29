"""Extension-dispatched plain-text extraction for the knowledge index.

Rich-doc parsers (pymupdf/python-docx/openpyxl/python-pptx) are imported
lazily so a missing optional dependency degrades to skipping that file
(returns None) rather than crashing the whole index rebuild.
"""

from __future__ import annotations

from pathlib import Path

from loguru import logger

_PLAIN_EXTS = {".md", ".txt", ".rst", ".json", ".yaml", ".yml", ".py"}


def _read_plain(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="replace")


def _read_pdf(path: Path) -> str:
    import fitz  # pymupdf

    parts: list[str] = []
    with fitz.open(path) as doc:
        for page in doc:
            parts.append(page.get_text())
    return "\n".join(parts)


def _read_docx(path: Path) -> str:
    import docx

    document = docx.Document(str(path))
    parts = [p.text for p in document.paragraphs]
    for table in document.tables:
        for row in table.rows:
            parts.append("\t".join(cell.text for cell in row.cells))
    return "\n".join(parts)


def _read_xlsx(path: Path) -> str:
    import openpyxl

    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    parts: list[str] = []
    for ws in wb.worksheets:
        parts.append(f"# {ws.title}")
        for row in ws.iter_rows(values_only=True):
            cells = ["" if v is None else str(v) for v in row]
            parts.append("\t".join(cells))
    wb.close()
    return "\n".join(parts)


def _read_pptx(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(shape.text_frame.text)
    return "\n".join(parts)


_RICH = {
    ".pdf": _read_pdf,
    ".docx": _read_docx,
    ".xlsx": _read_xlsx,
    ".pptx": _read_pptx,
}


def extract_text(path: Path) -> str | None:
    """Return extracted plain text, or None if the file should be skipped.

    Missing parser library or any parse error is logged at warning level and
    yields None so the caller skips the file without aborting the rebuild.
    """
    suffix = path.suffix.lower()
    reader = _RICH.get(suffix)
    try:
        if reader is not None:
            return reader(path)
        return _read_plain(path)
    except ImportError as e:
        logger.warning("Parser library missing for {} ({}); skipping", path, e)
        return None
    except Exception as e:
        logger.warning("Failed to extract text from {}: {}", path, e)
        return None
